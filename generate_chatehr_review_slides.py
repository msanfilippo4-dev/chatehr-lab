"""
ChatEHR Lab Review: What You Discovered
Generates chatehr_review.pptx (~28 slides) — a debrief deck covering all 3
directed patient cases, student highlights, and key takeaways.
"""

import os
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Rectangle
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parents[1]
LAB_DIR = Path(__file__).resolve().parent
IMAGES_DIR = LAB_DIR / "images_chatehr_review"
TEMPLATE = ROOT / "powerpoint_theme" / "MASTER BASIC PPT DECK TEMPLATE.pptx"
OUTPUT = LAB_DIR / "chatehr_review.pptx"

# ---------------------------------------------------------------------------
# Brand colors
# ---------------------------------------------------------------------------
FORDHAM_MAROON = RGBColor(139, 21, 56)
FORDHAM_DARK = RGBColor(51, 51, 51)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(212, 175, 55)
SURVEY_GREEN = RGBColor(39, 174, 96)
QUIZ_BLUE = RGBColor(52, 152, 219)
BREAK_ORANGE = RGBColor(230, 126, 34)

FONT_NAME = "Arial"

# Layout constants for 10x7.5 design space
LEFT = Inches(0.5)
SECTION_TOP = Inches(0.3)
TITLE_TOP = Inches(0.6)
SUBTITLE_TOP = Inches(1.3)
BODY_TOP = Inches(1.8)
BODY_WIDTH = Inches(9.0)
BODY_HEIGHT = Inches(5.0)

# Matplotlib palette
C_MAROON = "#8B1538"
C_GOLD = "#D4AF37"
C_BLUE = "#3498DB"
C_GREEN = "#27AE60"
C_RED = "#E74C3C"
C_DARK = "#2E3138"
C_MUTED = "#7F8C8D"
C_MAROON_SOFT = "#F8E7ED"
C_BLUE_SOFT = "#EBF5FB"
C_GREEN_SOFT = "#EAFAF1"
C_GOLD_SOFT = "#FDF2D0"
C_RED_SOFT = "#FDEDEC"
C_GRAY_SOFT = "#EEF1F5"
C_PURPLE_SOFT = "#F4ECF7"
C_ORANGE_SOFT = "#FEF5E7"


# ---------------------------------------------------------------------------
# Visual generation helpers
# ---------------------------------------------------------------------------

def _new_canvas(title, subtitle=""):
    fig, ax = plt.subplots(figsize=(12.8, 7.2), dpi=150)
    fig.patch.set_facecolor("white")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.axis("off")
    ax.text(50, 95, title, ha="center", va="center", fontsize=22,
            color=C_MAROON, fontweight="bold")
    if subtitle:
        ax.text(50, 89, subtitle, ha="center", va="center", fontsize=12,
                color=C_DARK)
    return fig, ax


def _box(ax, x, y, w, h, text, fc="#F8F8F8", ec=C_MAROON, fs=11, tc=C_DARK):
    patch = FancyBboxPatch(
        (x, y), w, h,
        boxstyle="round,pad=0.5,rounding_size=2.0",
        linewidth=1.8, edgecolor=ec, facecolor=fc,
    )
    ax.add_patch(patch)
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center",
            fontsize=fs, color=tc, wrap=True)


def _arrow(ax, x1, y1, x2, y2, label=""):
    arr = FancyArrowPatch((x1, y1), (x2, y2), arrowstyle="-|>",
                          mutation_scale=14, linewidth=1.8, color=C_MAROON)
    ax.add_patch(arr)
    if label:
        ax.text((x1 + x2) / 2, (y1 + y2) / 2 + 2, label, ha="center",
                va="center", fontsize=10, color=C_DARK)


def _save_both(fig, stem):
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)
    svg_path = IMAGES_DIR / f"{stem}.svg"
    png_path = IMAGES_DIR / f"{stem}.png"
    fig.savefig(str(svg_path), format="svg", bbox_inches="tight")
    fig.savefig(str(png_path), format="png", dpi=190, bbox_inches="tight")
    plt.close(fig)
    return {"svg": str(svg_path), "png": str(png_path)}


# ---------------------------------------------------------------------------
# 6 Visual Figures
# ---------------------------------------------------------------------------

def generate_visuals():
    visuals = {}

    # 01 — Score Distribution
    fig, ax = plt.subplots(figsize=(12.8, 7.2), dpi=150)
    fig.patch.set_facecolor("white")
    students = ["Brenda", "James", "Matias", "NYR2", "Natalie",
                "Penelope", "Student F", "Student G", "Student H", "Student J"]
    best_scores = [100, 95, 95, 93, 91, 89, 89, 88, 85, 82]
    colors = []
    for s in best_scores:
        if s >= 95:
            colors.append(C_GREEN)
        elif s >= 90:
            colors.append(C_BLUE)
        elif s >= 85:
            colors.append(C_GOLD)
        else:
            colors.append(C_MUTED)
    bars = ax.bar(range(len(students)), best_scores, color=colors, width=0.6,
                  edgecolor="white", linewidth=1.5, zorder=3)
    ax.set_xticks(range(len(students)))
    ax.set_xticklabels(students, fontsize=11, color=C_DARK, rotation=30, ha="right")
    ax.set_ylabel("Best Score (0-100)", fontsize=12, color=C_DARK)
    ax.set_ylim(0, 110)
    ax.set_title("Class Score Distribution — Best Score per Student",
                 fontsize=18, color=C_MAROON, fontweight="bold", pad=15)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for i, v in enumerate(best_scores):
        ax.text(i, v + 2, str(v), ha="center", fontsize=12,
                fontweight="bold", color=C_DARK)
    legend_patches = [
        mpatches.Patch(color=C_GREEN, label="95-100"),
        mpatches.Patch(color=C_BLUE, label="90-94"),
        mpatches.Patch(color=C_GOLD, label="85-89"),
        mpatches.Patch(color=C_MUTED, label="<85"),
    ]
    ax.legend(handles=legend_patches, fontsize=10, loc="upper right",
              title="Score Range")
    plt.tight_layout()
    visuals["score_distribution"] = _save_both(fig, "01_score_distribution")

    # 02 — RAG Comparison
    fig, ax = _new_canvas("RAG OFF vs RAG ON",
                          "How retrieval-augmented generation changes clinical assistant behavior")
    # RAG OFF column
    _box(ax, 3, 52, 44, 30, "", fc=C_RED_SOFT, ec=C_RED)
    ax.text(25, 78, "RAG OFF", ha="center", fontsize=16,
            color=C_RED, fontweight="bold")
    ax.text(25, 72, "Generic Guidelines Mode", ha="center", fontsize=11,
            color=C_DARK)
    items_off = [
        '"Consider checking HbA1c..."',
        '"If lipids are elevated, may need..."',
        '"General guidelines suggest..."',
        'No patient-specific values cited',
        'Hedging language throughout',
    ]
    for i, item in enumerate(items_off):
        ax.text(25, 64 - i * 6, item, ha="center", fontsize=9.5, color=C_DARK)

    # RAG ON column
    _box(ax, 53, 52, 44, 30, "", fc=C_GREEN_SOFT, ec=C_GREEN)
    ax.text(75, 78, "RAG ON", ha="center", fontsize=16,
            color=C_GREEN, fontweight="bold")
    ax.text(75, 72, "Chart-Grounded Mode", ha="center", fontsize=11,
            color=C_DARK)
    items_on = [
        '"HbA1c is 8.7% — above ADA target"',
        '"LDL 142 mg/dL — add high-intensity statin"',
        '"Recommend GLP-1 RA (note: sulfa allergy)"',
        'Specific values, dates, medications',
        'Confident, evidence-grounded language',
    ]
    for i, item in enumerate(items_on):
        ax.text(75, 64 - i * 6, item, ha="center", fontsize=9.5, color=C_DARK)

    ax.text(50, 22, "Key insight: RAG ON improves specificity but can create "
            "false confidence from incomplete retrieval",
            ha="center", fontsize=11, color=C_DARK, fontweight="bold")
    visuals["rag_comparison"] = _save_both(fig, "02_rag_comparison")

    # 03 — Context Levels
    fig, ax = _new_canvas("Three Context Levels — Minimum Necessary Principle",
                          "Chloe Davis (LAB-003): How much should the AI see?")
    # LIMITED
    _box(ax, 3, 55, 28, 26, "", fc=C_BLUE_SOFT, ec=C_BLUE)
    ax.text(17, 78, "LIMITED", ha="center", fontsize=14,
            color=C_BLUE, fontweight="bold")
    limited_items = ["Age, gender, zip prefix", "No conditions or meds",
                     "No lab results", "No visit notes"]
    for i, item in enumerate(limited_items):
        ax.text(17, 72 - i * 5, item, ha="center", fontsize=9, color=C_DARK)

    _arrow(ax, 31, 68, 36, 68, "")

    # STANDARD
    _box(ax, 36, 55, 28, 26, "", fc=C_GREEN_SOFT, ec=C_GREEN)
    ax.text(50, 78, "STANDARD", ha="center", fontsize=14,
            color=C_GREEN, fontweight="bold")
    standard_items = ["Demographics + conditions", "All labs (hCG 450)",
                      "All meds (Isotretinoin)", "Allergies & immunizations"]
    for i, item in enumerate(standard_items):
        ax.text(50, 72 - i * 5, item, ha="center", fontsize=9, color=C_DARK)

    _arrow(ax, 64, 68, 69, 68, "")

    # FULL
    _box(ax, 69, 55, 28, 26, "", fc=C_MAROON_SOFT, ec=C_MAROON)
    ax.text(83, 78, "FULL", ha="center", fontsize=14,
            color=C_MAROON, fontweight="bold")
    full_items = ["Everything in STANDARD", "+ Visit history & notes",
                  "+ Provider already acted", "+ Stop isotretinoin order"]
    for i, item in enumerate(full_items):
        ax.text(83, 72 - i * 5, item, ha="center", fontsize=9, color=C_DARK)

    # Bottom note
    _box(ax, 10, 22, 80, 16,
         "Privacy exposure increases left to right.\n"
         "Clinical usefulness also increases — but is FULL always necessary?\n"
         "STANDARD reveals the teratogen + pregnancy without exposing provider actions.",
         fc=C_GRAY_SOFT, ec=C_MUTED, fs=10)
    visuals["context_levels"] = _save_both(fig, "03_context_levels")

    # 04 — Eleanor Vance Case Card
    fig, ax = _new_canvas("Patient Summary: Eleanor Vance (LAB-001)",
                          "71F — Heart Failure, Hypertension")
    _box(ax, 5, 52, 40, 30, "", fc=C_MAROON_SOFT, ec=C_MAROON)
    ax.text(25, 78, "Key Labs", ha="center", fontsize=14,
            color=C_MAROON, fontweight="bold")
    labs_e = [
        "Creatinine: 1.4 mg/dL (High)",
        "BNP: 350 pg/mL (High)",
        "Potassium: 4.2 mEq/L (Normal — March 1)",
        "K+ 5.8 mEq/L (in visit notes — March 15)",
    ]
    for i, lab in enumerate(labs_e):
        color = C_RED if "5.8" in lab else C_DARK
        weight = "bold" if "5.8" in lab else "normal"
        ax.text(25, 72 - i * 6, lab, ha="center", fontsize=10,
                color=color, fontweight=weight)

    _box(ax, 55, 52, 40, 30, "", fc=C_BLUE_SOFT, ec=C_BLUE)
    ax.text(75, 78, "Medications", ha="center", fontsize=14,
            color=C_BLUE, fontweight="bold")
    meds_e = [
        "Carvedilol 12.5mg BID (Active)",
        "Lisinopril 20mg Daily (Active)",
        "Furosemide 40mg Daily (Active)",
        "Spironolactone 25mg (DISCONTINUED)",
    ]
    for i, med in enumerate(meds_e):
        color = C_RED if "DISC" in med else C_DARK
        weight = "bold" if "DISC" in med else "normal"
        ax.text(75, 72 - i * 6, med, ha="center", fontsize=10,
                color=color, fontweight=weight)

    _box(ax, 15, 18, 70, 16,
         "RISK: NSAIDs create 'Triple Whammy' — ACE inhibitor + Diuretic + "
         "NSAID = Acute Kidney Injury\n"
         "Spironolactone restart dangerous: K+ already 5.8 + ACE inhibitor",
         fc=C_RED_SOFT, ec=C_RED, fs=10)
    visuals["eleanor_case"] = _save_both(fig, "04_eleanor_case")

    # 05 — Marcus Thorne Case Card
    fig, ax = _new_canvas("Patient Summary: Marcus Thorne (LAB-002)",
                          "57M — Type 2 Diabetes, Hyperlipidemia")
    _box(ax, 5, 52, 40, 30, "", fc=C_MAROON_SOFT, ec=C_MAROON)
    ax.text(25, 78, "Key Labs", ha="center", fontsize=14,
            color=C_MAROON, fontweight="bold")
    labs_m = [
        "Fasting Glucose: 165 mg/dL (High)",
        "Total Cholesterol: 210 mg/dL (High)",
        "Creatinine: 0.9 mg/dL (Normal)",
        "HbA1c: 8.7% (in visit notes only)",
        "LDL: 142 mg/dL (in visit notes only)",
    ]
    for i, lab in enumerate(labs_m):
        color = C_RED if "notes" in lab else C_DARK
        weight = "bold" if "notes" in lab else "normal"
        ax.text(25, 72 - i * 5.5, lab, ha="center", fontsize=10,
                color=color, fontweight=weight)

    _box(ax, 55, 52, 40, 30, "", fc=C_BLUE_SOFT, ec=C_BLUE)
    ax.text(75, 78, "Medications & Allergy", ha="center", fontsize=14,
            color=C_BLUE, fontweight="bold")
    meds_m = [
        "Metformin 1000mg BID (Active)",
        "Atorvastatin 40mg Nightly (Active)",
        "",
        "ALLERGY: Sulfa drugs (Hives)",
    ]
    for i, med in enumerate(meds_m):
        color = C_RED if "ALLERGY" in med else C_DARK
        weight = "bold" if "ALLERGY" in med else "normal"
        ax.text(75, 72 - i * 6, med, ha="center", fontsize=10,
                color=color, fontweight=weight)

    _box(ax, 15, 18, 70, 16,
         "KEY INSIGHT: HbA1c 8.7% and LDL 142 are NOT in the structured labs.\n"
         "They appear only in the visit note — from the cardiologist's "
         "point-of-care labs.\nRAG ON surfaces these; RAG OFF cannot.",
         fc=C_GOLD_SOFT, ec=C_GOLD, fs=10)
    visuals["marcus_case"] = _save_both(fig, "05_marcus_case")

    # 06 — Chloe Davis Case Card
    fig, ax = _new_canvas("Patient Summary: Chloe Davis (LAB-003)",
                          "27F — Acne (Isotretinoin), Positive Pregnancy Test")
    _box(ax, 5, 52, 40, 30, "", fc=C_MAROON_SOFT, ec=C_MAROON)
    ax.text(25, 78, "Key Labs", ha="center", fontsize=14,
            color=C_MAROON, fontweight="bold")
    labs_c = [
        "Pregnancy test (hCG): 450 mIU/mL",
        "Hemoglobin: 13.5 g/dL (Normal)",
    ]
    for i, lab in enumerate(labs_c):
        color = C_RED if "hCG" in lab else C_DARK
        weight = "bold" if "hCG" in lab else "normal"
        ax.text(25, 72 - i * 6, lab, ha="center", fontsize=10,
                color=color, fontweight=weight)

    _box(ax, 55, 52, 40, 30, "", fc=C_BLUE_SOFT, ec=C_BLUE)
    ax.text(75, 78, "Medications", ha="center", fontsize=14,
            color=C_BLUE, fontweight="bold")
    meds_c = [
        "Isotretinoin 40mg Daily (ACTIVE)",
        "  — Category X teratogen",
        "Ethinyl estradiol/Norethindrone",
        "Albuterol inhaler PRN",
    ]
    for i, med in enumerate(meds_c):
        color = C_RED if "Category X" in med or "Isotretinoin" in med else C_DARK
        weight = "bold" if "Category X" in med or "Isotretinoin" in med else "normal"
        ax.text(75, 72 - i * 6, med, ha="center", fontsize=10,
                color=color, fontweight=weight)

    _box(ax, 15, 18, 70, 16,
         "CRITICAL: Isotretinoin is a Category X teratogen — must stop "
         "immediately.\nVisit notes reveal provider already acted (April 2 "
         "telehealth).\nPrivacy question: How much context does the AI need "
         "to flag this safely?",
         fc=C_RED_SOFT, ec=C_RED, fs=10)
    visuals["chloe_case"] = _save_both(fig, "06_chloe_case")

    return visuals


# =========================================================================
# FordhamSlideFactory  (copied from week7 — canonical pattern)
# =========================================================================

class FordhamSlideFactory:

    def __init__(self, template_path):
        self.prs = Presentation(str(template_path))
        self.design_w = Inches(10)
        self.design_h = Inches(7.5)
        self.scale_x = self.prs.slide_width / self.design_w
        self.scale_y = self.prs.slide_height / self.design_h
        self.font_scale = min(self.scale_x, self.scale_y)
        self._clear_slides()

    def _sx(self, val):
        return int(val * self.scale_x)

    def _sy(self, val):
        return int(val * self.scale_y)

    def _sp(self, pt):
        return max(11, int(round(pt * self.font_scale)))

    def _clear_slides(self):
        slide_id_list = self.prs.slides._sldIdLst
        for sld_id in list(slide_id_list):
            r_id = sld_id.rId
            self.prs.part.drop_rel(r_id)
            slide_id_list.remove(sld_id)

    def _set_font(self, p, size, color=FORDHAM_DARK, bold=False, italic=False,
                  align=PP_ALIGN.LEFT):
        scaled = Pt(self._sp(size))
        p.font.name = FONT_NAME
        p.font.size = scaled
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.alignment = align
        for run in p.runs:
            run.font.name = FONT_NAME
            run.font.size = scaled
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.italic = italic

    def _add_section_header(self, slide, text):
        box = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(SECTION_TOP),
            self._sx(Inches(4)), self._sy(Inches(0.25)))
        p = box.text_frame.paragraphs[0]
        p.text = text.upper()
        self._set_font(p, 10, color=FORDHAM_MAROON)

    def _add_title(self, slide, text, size=28):
        box = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(TITLE_TOP),
            self._sx(BODY_WIDTH), self._sy(Inches(0.6)))
        p = box.text_frame.paragraphs[0]
        p.text = text
        self._set_font(p, size, color=FORDHAM_MAROON, bold=True)

    def _add_subtitle(self, slide, text):
        box = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(SUBTITLE_TOP),
            self._sx(BODY_WIDTH), self._sy(Inches(0.4)))
        p = box.text_frame.paragraphs[0]
        p.text = text
        self._set_font(p, 14, color=FORDHAM_MAROON)

    def _choose_bullet_size(self, items):
        clean = [x.strip() for x in items if x and x.strip()]
        count = len(clean)
        longest = max((len(x) for x in clean), default=0)
        if count <= 4 and longest <= 72:
            return 20
        if count <= 5 and longest <= 90:
            return 18
        if count <= 6:
            return 17
        return 16

    def _add_bullets(self, slide, items, top=BODY_TOP):
        box = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(top),
            self._sx(BODY_WIDTH), self._sy(BODY_HEIGHT))
        tf = box.text_frame
        tf.word_wrap = True
        size = self._choose_bullet_size(items)
        first = True
        for item in items:
            if not item or not str(item).strip():
                continue
            text = str(item).strip()
            is_heading = text.startswith("H:")
            if is_heading:
                text = text[2:].strip()
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            if is_heading:
                p.text = text
                self._set_font(p, size + 1, color=FORDHAM_DARK, bold=True)
            else:
                p.text = f"• {text}"
                self._set_font(p, size, color=FORDHAM_DARK)
            p.space_after = Pt(max(2, int(round(3 * self.font_scale))))

    def _add_notes(self, slide, notes):
        if notes and notes.strip():
            slide.notes_slide.notes_text_frame.text = notes.strip()

    def _set_background(self, slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    # --- Public slide builders ---

    def add_title_slide(self, title, subtitle, notes):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        self._set_background(slide, WHITE)
        tb = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(Inches(2.5)),
            self._sx(BODY_WIDTH), self._sy(Inches(1.0)))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        self._set_font(p, 36, color=FORDHAM_MAROON, bold=True, align=PP_ALIGN.CENTER)
        sb = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(Inches(3.8)),
            self._sx(BODY_WIDTH), self._sy(Inches(1.5)))
        sb.text_frame.word_wrap = True
        for i, line in enumerate(subtitle.split('\n')):
            pp = sb.text_frame.paragraphs[0] if i == 0 else sb.text_frame.add_paragraph()
            pp.text = line
            self._set_font(pp, 16, color=FORDHAM_MAROON, align=PP_ALIGN.CENTER)
        self._add_notes(slide, notes)
        return slide

    def add_section_divider(self, part, subtitle, notes):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        self._set_background(slide, FORDHAM_MAROON)
        box = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(Inches(2.8)),
            self._sx(BODY_WIDTH), self._sy(Inches(0.8)))
        p = box.text_frame.paragraphs[0]
        p.text = part
        self._set_font(p, 36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        box2 = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(Inches(3.8)),
            self._sx(BODY_WIDTH), self._sy(Inches(0.6)))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = subtitle
        self._set_font(p2, 18, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
        self._add_notes(slide, notes)
        return slide

    def add_content_slide(self, title, bullets, notes, section=None, subtitle=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        if section:
            self._add_section_header(slide, section)
        self._add_title(slide, title)
        if subtitle:
            self._add_subtitle(slide, subtitle)
        self._add_bullets(slide, bullets)
        self._add_notes(slide, notes)
        return slide

    def add_image_slide(self, title, image_path, notes, section=None, caption=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        if section:
            self._add_section_header(slide, section)
        self._add_title(slide, title, size=24)
        img_path = str(image_path)
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(
                    img_path,
                    self._sx(Inches(0.4)), self._sy(Inches(1.15)),
                    width=self._sx(Inches(9.2)), height=self._sy(Inches(5.75)))
            except Exception as exc:
                box = slide.shapes.add_textbox(
                    self._sx(LEFT), self._sy(Inches(3)),
                    self._sx(BODY_WIDTH), self._sy(Inches(1)))
                p = box.text_frame.paragraphs[0]
                p.text = f"[Image: {os.path.basename(img_path)}] ({exc})"
                self._set_font(p, 14, color=FORDHAM_MAROON, italic=True,
                               align=PP_ALIGN.CENTER)
        else:
            box = slide.shapes.add_textbox(
                self._sx(LEFT), self._sy(Inches(3)),
                self._sx(BODY_WIDTH), self._sy(Inches(1)))
            p = box.text_frame.paragraphs[0]
            p.text = f"[Missing: {os.path.basename(img_path)}]"
            self._set_font(p, 14, color=FORDHAM_MAROON, italic=True,
                           align=PP_ALIGN.CENTER)
        if caption:
            cb = slide.shapes.add_textbox(
                self._sx(LEFT), self._sy(Inches(6.8)),
                self._sx(BODY_WIDTH), self._sy(Inches(0.3)))
            p = cb.text_frame.paragraphs[0]
            p.text = caption
            self._set_font(p, 10, color=FORDHAM_DARK, italic=True,
                           align=PP_ALIGN.CENTER)
        self._add_notes(slide, notes)
        return slide

    def add_table_slide(self, title, headers, rows, notes, section=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        if section:
            self._add_section_header(slide, section)
        self._add_title(slide, title)
        num_cols, num_rows = len(headers), len(rows) + 1
        tbl = slide.shapes.add_table(
            num_rows, num_cols,
            self._sx(LEFT), self._sy(BODY_TOP),
            self._sx(BODY_WIDTH), self._sy(Inches(4.0))).table
        hsize = 12 if num_rows <= 5 else 11 if num_rows <= 7 else 10
        bsize = hsize - 1
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = FORDHAM_MAROON
            self._set_font(cell.text_frame.paragraphs[0], hsize, color=WHITE,
                           bold=True)
        for r, row in enumerate(rows, 1):
            for c, val in enumerate(row):
                tbl.cell(r, c).text = str(val)
                self._set_font(tbl.cell(r, c).text_frame.paragraphs[0], bsize)
        self._add_notes(slide, notes)
        return slide

    def add_two_column_slide(self, title, left_title, left_items,
                             right_title, right_items, notes, section=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        if section:
            self._add_section_header(slide, section)
        self._add_title(slide, title)
        ltb = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(SUBTITLE_TOP),
            self._sx(Inches(4.2)), self._sy(Inches(0.4)))
        self._set_font(ltb.text_frame.paragraphs[0], 14,
                       color=FORDHAM_MAROON, bold=True)
        ltb.text_frame.paragraphs[0].text = left_title
        lbox = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(Inches(1.8)),
            self._sx(Inches(4.2)), self._sy(Inches(4.5)))
        lbox.text_frame.word_wrap = True
        first = True
        for item in left_items:
            if not item:
                continue
            p = lbox.text_frame.paragraphs[0] if first else lbox.text_frame.add_paragraph()
            first = False
            p.text = f"• {item}"
            self._set_font(p, 20)
        rx = Inches(5.0)
        rtb = slide.shapes.add_textbox(
            self._sx(rx), self._sy(SUBTITLE_TOP),
            self._sx(Inches(4.2)), self._sy(Inches(0.4)))
        self._set_font(rtb.text_frame.paragraphs[0], 14,
                       color=FORDHAM_MAROON, bold=True)
        rtb.text_frame.paragraphs[0].text = right_title
        rbox = slide.shapes.add_textbox(
            self._sx(rx), self._sy(Inches(1.8)),
            self._sx(Inches(4.2)), self._sy(Inches(4.5)))
        rbox.text_frame.word_wrap = True
        first = True
        for item in right_items:
            if not item:
                continue
            p = rbox.text_frame.paragraphs[0] if first else rbox.text_frame.add_paragraph()
            first = False
            p.text = f"• {item}"
            self._set_font(p, 20)
        self._add_notes(slide, notes)
        return slide

    def add_quote_slide(self, quote, attribution, notes, section=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        if section:
            self._add_section_header(slide, section)
        qb = slide.shapes.add_textbox(
            self._sx(Inches(0.8)), self._sy(Inches(2.0)),
            self._sx(Inches(8.4)), self._sy(Inches(3.0)))
        qb.text_frame.word_wrap = True
        p = qb.text_frame.paragraphs[0]
        p.text = f'"{quote}"'
        self._set_font(p, 20, color=FORDHAM_DARK, italic=True, align=PP_ALIGN.CENTER)
        ab = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(Inches(5.5)),
            self._sx(BODY_WIDTH), self._sy(Inches(0.4)))
        ab.text_frame.paragraphs[0].text = f"— {attribution}"
        self._set_font(ab.text_frame.paragraphs[0], 14, color=FORDHAM_MAROON,
                       bold=True, align=PP_ALIGN.CENTER)
        self._add_notes(slide, notes)
        return slide

    def add_closing_slide(self, title, subtitle, notes):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        self._set_background(slide, FORDHAM_MAROON)
        tb = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(Inches(2.8)),
            self._sx(BODY_WIDTH), self._sy(Inches(0.8)))
        tb.text_frame.paragraphs[0].text = title
        self._set_font(tb.text_frame.paragraphs[0], 36, color=WHITE, bold=True,
                       align=PP_ALIGN.CENTER)
        sb = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(Inches(3.8)),
            self._sx(BODY_WIDTH), self._sy(Inches(0.6)))
        sb.text_frame.paragraphs[0].text = subtitle
        self._set_font(sb.text_frame.paragraphs[0], 18, color=WHITE,
                       align=PP_ALIGN.CENTER)
        ib = slide.shapes.add_textbox(
            self._sx(LEFT), self._sy(Inches(4.8)),
            self._sx(BODY_WIDTH), self._sy(Inches(0.6)))
        ib.text_frame.paragraphs[0].text = (
            "HINF 6117 | ChatEHR Lab Review | Fordham University")
        self._set_font(ib.text_frame.paragraphs[0], 12, color=WHITE,
                       align=PP_ALIGN.CENTER)
        self._add_notes(slide, notes)
        return slide

    def save(self, output_path):
        self.prs.save(str(output_path))
        print(f"Saved: {output_path} ({len(self.prs.slides)} slides)")


# =========================================================================
# Slide builders
# =========================================================================

def build_opening(f, visuals):
    """Slides 1-4: Title, What We Did, Score Distribution, Participation."""

    # 1 — Title
    f.add_title_slide(
        "CHATEHR LAB REVIEW",
        "What You Discovered\nHINF 6117 | Fordham University | Spring 2026",
        "Welcome back. Last week you got your hands on a simulated EHR with a "
        "Gemini-powered clinical assistant, and I've now read through all 23 of your "
        "submissions. Today we're going to debrief — go over the correct answers for "
        "each case, highlight the strongest responses, and lock in the key lessons "
        "about AI-assisted clinical decision-making. Let's get into it."
    )

    # 2 — What We Did
    f.add_content_slide("What We Did", [
        "H:Three Missions — 60 Minutes",
        "Mission A: Eleanor Vance (LAB-001) — Medication safety review",
        "Mission B: Marcus Thorne (LAB-002) — RAG ON vs RAG OFF comparison",
        "Mission C: Chloe Davis (LAB-003) — Minimum necessary context decision",
        "",
        "H:Team Structure",
        "Navigator + Driver + Scribe — rotating roles across missions",
        "Each team submitted a reflection with three scored questions (Q1/Q2/Q3)",
    ], "Here's what you did last week. Three patients, three missions, 60 minutes. "
       "Each mission tested something different. Eleanor was about catching an unsafe "
       "medication recommendation. Marcus was about seeing the difference RAG makes. "
       "And Chloe was about the privacy tradeoff — how much context should the AI "
       "actually see? You worked in teams with rotating roles, and every team "
       "submitted a scored reflection. Let's see how you did.",
       section="OPENING")

    # 3 — Score Distribution
    f.add_image_slide(
        "Class Score Distribution",
        visuals["score_distribution"]["png"],
        "So here's the big picture. Most of you scored in the high 80s to 100 range, "
        "which tells me you understood the core concepts. Brenda's team hit 100 on "
        "their best submission. A few of you submitted multiple times as you refined "
        "your answers, which is exactly what I wanted — that iterative improvement "
        "is part of the learning. The scores here show your single best submission.",
        section="OPENING",
        caption="Best score per student across all submissions"
    )

    # 4 — Participation Summary
    f.add_content_slide("Participation Summary", [
        "10 unique students submitted reflections",
        "23 total submissions (some students refined and resubmitted)",
        "3 Health Bench advanced task submissions",
        "1 Extra Credit submission (Brenda — Bronx population analysis)",
        "Average best score: ~91 across all students",
        "Most common strength: identifying medication risks with chart evidence",
    ], "Ten of you submitted, and several submitted more than once. That's 23 total "
       "reflections I read through. Three of you attempted the Health Bench — that's "
       "the advanced task where you evaluated the AI's clinical reasoning against a "
       "rubric. And Brenda went further with the extra credit, doing a population-level "
       "analysis on the Bronx patient panel. I'm going to highlight specific responses "
       "as we go through each case.",
       section="OPENING")


def build_q1_review(f, visuals):
    """Slides 5-10: Eleanor Vance — Medication Safety."""

    # 5 — Section Divider
    f.add_section_divider(
        "Mission A: Eleanor Vance",
        "Medication Safety Review",
        "Alright, let's start with Eleanor. This was the medication safety case. "
        "Your job was to ask the AI assistant about pain management for a heart "
        "failure patient and then evaluate whether its recommendation was safe."
    )

    # 6 — The Case
    f.add_image_slide(
        "Eleanor Vance — The Case",
        visuals["eleanor_case"]["png"],
        "Here's Eleanor at a glance. 71-year-old woman with heart failure and "
        "hypertension. Her structured labs from March 1 show a creatinine of 1.4 — "
        "that's already borderline — a BNP of 350, and a potassium of 4.2 which "
        "looks normal. She's on Carvedilol, Lisinopril, and Furosemide — that's a "
        "beta blocker, an ACE inhibitor, and a loop diuretic. And notice that "
        "Spironolactone is listed as discontinued. Hold onto that — it matters.",
        section="Q1 — ELEANOR VANCE"
    )

    # 7 — The Unsafe Recommendation
    f.add_two_column_slide(
        "The Unsafe Recommendation",
        "UNSAFE: NSAIDs",
        [
            "Creates the 'Triple Whammy'",
            "ACE inhibitor + Diuretic + NSAID",
            "Reduces renal blood flow",
            "Risk of acute kidney injury",
            "Creatinine already elevated (1.4)",
        ],
        "SAFER Alternatives",
        [
            "Acetaminophen (Tylenol) first-line",
            "Topical analgesics (lidocaine patches)",
            "Non-pharmacological approaches",
            "If NSAID needed: lowest dose, shortest duration",
            "Monitor renal function closely",
        ],
        "This is the core teaching point. If the AI recommended NSAIDs — ibuprofen, "
        "naproxen, anything in that class — that's a safety issue. Eleanor is already "
        "on an ACE inhibitor and a diuretic. You add an NSAID to that combination "
        "and you get what's called the 'Triple Whammy' — three drugs that all reduce "
        "renal perfusion in different ways. With her creatinine already at 1.4, that's "
        "a real risk for acute kidney injury. The safe answer is acetaminophen as "
        "first-line, or topical options.",
        section="Q1 — ELEANOR VANCE"
    )

    # 8 — Best Answers
    f.add_content_slide("Your Best Answers", [
        'H:NYR2 / Group 3',
        '"Identified the triple whammy — ACE inhibitor + diuretic + NSAID — '
        'and cited BNP 350 as evidence of active heart failure"',
        'H:Penelope & Natalie',
        '"Flagged the Spironolactone restart risk — K+ was 5.8 with an ACE '
        'inhibitor already on board, making hyperkalemia dangerous"',
        'H:Matias',
        '"Do not take potassium supplements — K+ already 5.8 mEq/L. '
        'Clear, specific, grounded in the chart"',
    ], "Let me highlight a few standout answers. NYR2's group nailed the triple "
       "whammy identification and didn't just name it — they cited the BNP of 350 "
       "as evidence that her heart failure is active, which makes the NSAID risk "
       "even worse. Penelope and Natalie caught something else: the spironolactone "
       "restart risk. Even though it's discontinued, some of your AI assistants "
       "suggested restarting it. But the visit notes reveal her potassium was 5.8 — "
       "that's dangerously high for someone on an ACE inhibitor. And Matias gave a "
       "really clean, specific warning about potassium supplements.",
       section="Q1 — ELEANOR VANCE")

    # 9 — The Hidden Lesson
    f.add_content_slide("The Hidden Lesson", [
        "H:Two Different Potassium Values — Two Different Locations",
        "Labs tab (structured data, March 1): K+ = 4.2 mEq/L — Normal",
        "Visit notes (narrative text, March 15): K+ = 5.8 mEq/L — Dangerous",
        "",
        "H:Why This Matters",
        "The structured labs look reassuring — potassium is normal",
        "But the visit note reveals outside labs from ~March 12 with K+ 5.8",
        "The plan and notes sections of the visit contain the real clinical story",
        "If you only checked the Labs tab, you'd miss the most important finding",
    ], "Here's the hidden lesson in this case, and it's one of the most important "
       "takeaways from the whole lab. The Labs tab shows potassium of 4.2 — that's "
       "from March 1, and it looks perfectly normal. But if you read the visit note "
       "from March 15 — the plan and notes sections — Dr. Chen mentions that outside "
       "labs drawn three days earlier showed a potassium of 5.8. That's a critically "
       "different number. The structured data and the narrative data tell different "
       "stories, and the newer, more dangerous value is buried in the note. This is "
       "exactly what happens in real EHRs. The AI assistant with RAG can surface this, "
       "but you as the clinician need to know where to look.",
       section="Q1 — ELEANOR VANCE")

    # 10 — Q1 Takeaway
    f.add_content_slide("Q1 Takeaway", [
        "Always cite specific chart evidence — don't just say 'there's a risk'",
        "Check visit notes, not just structured lab tabs — key values hide in narrative",
        "Active vs. discontinued status matters — discontinued doesn't mean safe to restart",
        "The 'Triple Whammy' (ACE + diuretic + NSAID) is a classic med safety pattern",
        "Elevated K+ (5.8) + ACE inhibitor = do NOT add spironolactone or K+ supplements",
    ], "So here are your takeaways from Eleanor. First, specificity matters — cite the "
       "actual values, dates, and medications. Second, the Labs tab doesn't have the "
       "full picture. The visit notes contain values that haven't been entered as "
       "structured data yet. Third, just because a medication is discontinued doesn't "
       "mean the AI or you should recommend restarting it without checking why it was "
       "stopped. And fourth, the triple whammy and the K+ 5.8 — these are the kind of "
       "patterns that a well-informed clinician catches. The AI can help surface them, "
       "but only if you're reading critically.",
       section="Q1 — ELEANOR VANCE")


def build_q2_review(f, visuals):
    """Slides 11-17: Marcus Thorne — RAG Impact."""

    # 11 — Section Divider
    f.add_section_divider(
        "Mission B: Marcus Thorne",
        "RAG ON vs RAG OFF",
        "Now let's talk about Marcus. This was the RAG comparison case. You ran the "
        "same clinical question with RAG turned off and then on, and compared the "
        "responses. This is where a lot of you had real 'aha' moments."
    )

    # 12 — The Case
    f.add_image_slide(
        "Marcus Thorne — The Case",
        visuals["marcus_case"]["png"],
        "Marcus is a 57-year-old man with type 2 diabetes and hyperlipidemia. His "
        "structured labs show fasting glucose of 165 — that's high — and total "
        "cholesterol of 210. Creatinine is normal at 0.9. He's on Metformin 1000mg "
        "twice daily and Atorvastatin 40mg. And he has a sulfa allergy. But here's "
        "the key: his HbA1c of 8.7% and LDL of 142 are NOT in the Labs tab. They "
        "only appear in the visit note, where the provider mentions they came from "
        "point-of-care labs done by his cardiologist.",
        section="Q2 — MARCUS THORNE"
    )

    # 13 — RAG OFF vs RAG ON
    f.add_image_slide(
        "RAG OFF vs RAG ON",
        visuals["rag_comparison"]["png"],
        "Here's the comparison side by side. With RAG off, the AI gives you textbook "
        "guidelines. 'Consider checking HbA1c.' 'If lipids are elevated, may need to '— "
        "it hedges constantly because it has no patient data to ground on. Everything "
        "is conditional. With RAG on, it can see Marcus's chart. It says, 'HbA1c is "
        "8.7 percent, above the ADA target of 7 percent. Recommend adding a GLP-1 "
        "receptor agonist.' It cites the LDL of 142 and suggests escalating statin "
        "therapy. Night and day difference. The language shifts from uncertain to "
        "confident, from generic to specific.",
        section="Q2 — MARCUS THORNE"
    )

    # 14 — What You Found
    f.add_content_slide("What You Found", [
        "H:Key Observations Across Submissions",
        "RAG ON responses were 2-3x more specific in medication recommendations",
        "RAG ON cited actual lab values (HbA1c 8.7%, LDL 142) vs. generic thresholds",
        "RAG OFF used hedging language: 'consider,' 'if indicated,' 'may need'",
        "RAG ON used confident language: 'recommend,' 'should be added,' 'is above target'",
        "Several teams noted RAG ON correctly flagged the sulfa allergy when recommending meds",
        "HbA1c and LDL were only available via visit notes — RAG ON surfaced them",
    ], "Across your submissions, a few themes jumped out. Everyone noticed the specificity "
       "increase — that's the obvious one. But several of you caught subtler things. The "
       "language shift from hedging to confident. The fact that RAG ON could pull HbA1c and "
       "LDL from the visit note even though they weren't in the structured labs. And some "
       "of you noticed that RAG ON correctly flagged the sulfa allergy when suggesting "
       "medications — that's the kind of safety check that only works when the AI can "
       "actually see the chart.",
       section="Q2 — MARCUS THORNE")

    # 15 — Student Highlight Quotes
    f.add_content_slide("Student Highlights", [
        'H:Brenda',
        '"RAG OFF is honest but clinically useless... every recommendation '
        'becomes a generic \'it depends\'"',
        'H:Matias',
        '"RAG OFF provides same result but grounded less on evidence... '
        'RAG ON more certain in language"',
        'H:James',
        'Provided a full structured comparison with evidence grounding checks '
        '— systematically evaluated each recommendation against chart data',
    ], "Three quotes that stuck with me. Brenda's is probably the most quotable — "
       "'RAG OFF is honest but clinically useless.' That's a perfect summary. The AI "
       "isn't wrong when RAG is off — it's just so generic that it can't help you make "
       "a decision for this specific patient. Matias noticed the grounding difference "
       "and the language certainty shift. And James did something really impressive — "
       "he built a structured comparison where he checked each recommendation against "
       "the actual chart evidence. That's the kind of systematic evaluation we want "
       "to see.",
       section="Q2 — MARCUS THORNE")

    # 16 — The Catch
    f.add_content_slide("The Catch: False Confidence", [
        "H:RAG ON Isn't Perfect",
        "Grounding does not equal correctness — the AI can confidently cite a value "
        "and still make a wrong recommendation",
        "Incomplete retrieval: if the RAG system misses a document, the AI won't "
        "know what it doesn't know",
        "Brenda's insight: 'RAG ON introduces false confidence from incomplete retrieval'",
        "",
        "H:The Clinician's Role",
        "RAG makes the AI more useful — but verification is still required",
        "Always cross-check AI recommendations against the full chart",
        "The AI's confidence level should not replace your own clinical judgment",
    ], "But here's the catch, and Brenda nailed it in her submission. RAG ON makes "
       "the AI sound more confident. It cites specific values, gives specific "
       "recommendations. But what if the retrieval missed something? What if there's "
       "a relevant lab or a contraindication that didn't get pulled into the context "
       "window? The AI won't say, 'Hey, I might be missing something.' It'll just "
       "give you a confident answer based on incomplete data. That's what Brenda calls "
       "'false confidence from incomplete retrieval.' And that's why the clinician still "
       "has to verify. RAG makes the AI more useful, but it doesn't make it infallible.",
       section="Q2 — MARCUS THORNE")

    # 17 — Q2 Takeaway
    f.add_content_slide("Q2 Takeaway", [
        "RAG dramatically improves specificity and clinical usefulness",
        "Without RAG, AI gives generic guidelines — with RAG, it gives patient-specific advice",
        "Key values (HbA1c 8.7%, LDL 142) can hide in narrative notes, not structured data",
        "Watch for overconfidence — grounded language doesn't guarantee correctness",
        "Clinician verification is non-negotiable, even with RAG ON",
    ], "So your takeaways for Marcus. RAG is a massive improvement — no question. But "
       "it's an improvement that comes with a new risk: overconfidence. The shift from "
       "'consider checking' to 'HbA1c is 8.7 percent' feels authoritative. And it is "
       "grounded in real data. But grounded doesn't mean complete, and complete doesn't "
       "mean correct. Your job as the clinician is to use that grounded information as "
       "a starting point, not an endpoint.",
       section="Q2 — MARCUS THORNE")


def build_q3_review(f, visuals):
    """Slides 18-23: Chloe Davis — Privacy Decision."""

    # 18 — Section Divider
    f.add_section_divider(
        "Mission C: Chloe Davis",
        "Minimum Necessary Context",
        "Our last case. Chloe was the privacy decision. This one had no single "
        "right answer — it was about reasoning through the tradeoff between clinical "
        "usefulness and privacy exposure."
    )

    # 19 — The Case
    f.add_image_slide(
        "Chloe Davis — The Case",
        visuals["chloe_case"]["png"],
        "Chloe is 27, being treated for severe acne with Isotretinoin — that's "
        "Accutane — which is a Category X teratogen. It causes severe birth defects. "
        "And she just had a positive pregnancy test — hCG of 450. She's also on birth "
        "control, so this was a contraceptive failure. The visit note from April 2 "
        "reveals that Dr. Patel already acted — counseled Chloe to stop isotretinoin "
        "immediately and referred her to high-risk OB. Your question was: how much "
        "context does the AI need to safely flag this situation?",
        section="Q3 — CHLOE DAVIS"
    )

    # 20 — Three Context Levels
    f.add_image_slide(
        "Three Context Levels",
        visuals["context_levels"]["png"],
        "Here are the three levels. LIMITED gives the AI almost nothing — age, gender, "
        "zip code. It can't identify any medication risk because it can't see the "
        "medication list. STANDARD gives it demographics, conditions, labs, and meds. "
        "So it can see the isotretinoin, it can see the positive hCG — that's enough "
        "to flag the teratogen risk. FULL adds the visit history and notes, so the AI "
        "also knows that the provider already acted on April 2. Each level reveals more, "
        "but also exposes more of the patient's private information to the AI system.",
        section="Q3 — CHLOE DAVIS"
    )

    # 21 — Your Choices
    f.add_content_slide("Your Choices", [
        "H:Most Teams Chose STANDARD",
        "Enough to flag Isotretinoin + pregnancy without exposing visit narrative",
        "Balanced clinical usefulness with privacy protection",
        "",
        "H:Some Teams Chose FULL",
        "Needed to see the April 2 note where the provider already acted",
        "Argument: without this context, the AI might duplicate an urgent intervention",
        "",
        "H:Strongest Reasoning: NYR2, Natalie, Penelope",
        'Explained why STANDARD over LIMITED and why STANDARD over FULL',
        "Didn't just pick a level — justified what each level adds and what it costs",
    ], "Most of you chose STANDARD, and that's a defensible answer. STANDARD gives the "
       "AI enough to see the isotretinoin and the positive pregnancy test — that's "
       "sufficient to flag the teratogen risk. But some of you chose FULL, and your "
       "reasoning was interesting. You argued that without the visit notes, the AI "
       "doesn't know the provider already acted. It might generate an urgent alert that's "
       "redundant. The strongest submissions — NYR2, Natalie, Penelope — didn't just pick "
       "a level. They explained what each level adds, what it costs in privacy exposure, "
       "and why their choice was the right balance.",
       section="Q3 — CHLOE DAVIS")

    # 22 — Student Highlight Quotes
    f.add_content_slide("Student Highlights", [
        'H:NYR2',
        '"STANDARD provides high clinical usefulness with moderate privacy risk"',
        'H:Matias',
        '"With both limited and standard, the EHR assistant cannot make a '
        'grounded response... With full context, it makes a specific response"',
        'H:James',
        '"Full context is necessary to understand the urgency of the '
        'intervention already taken"',
    ], "Three different perspectives here. NYR2 gave a clean cost-benefit framing — "
       "clinical usefulness versus privacy risk. That's the HIPAA minimum necessary "
       "principle in action. Matias made an interesting argument that even STANDARD "
       "wasn't enough for a truly grounded response — though I'd push back on that a "
       "bit, since STANDARD does include the medication list and pregnancy test. And "
       "James argued for FULL because he wanted the AI to know about the intervention "
       "that already happened. All three are reasonable — the point is that you "
       "reasoned through the tradeoff.",
       section="Q3 — CHLOE DAVIS")

    # 23 — Q3 Takeaway
    f.add_content_slide("Q3 Takeaway", [
        "'Minimum necessary' is context-dependent — there's no universal right answer",
        "More data does not always equal better care",
        "Privacy exposure should match clinical need — justify what you include AND exclude",
        "STANDARD catches the teratogen risk; FULL adds awareness of prior action",
        "The best answers explained the tradeoff, not just the choice",
    ], "Here's the takeaway. Minimum necessary isn't a fixed rule — it's a judgment "
       "call that depends on the clinical scenario. More data gives the AI more to work "
       "with, but it also exposes more of the patient's private information to the system. "
       "The best answers in this section didn't just pick LIMITED, STANDARD, or FULL. They "
       "explained what each level reveals, what it hides, and why the balance they chose "
       "was appropriate for this specific patient. That's the skill I want you to carry "
       "forward.",
       section="Q3 — CHLOE DAVIS")


def build_wrapup(f, visuals):
    """Slides 24-28: Health Bench, Patterns, Five Things, Debrief, Closing."""

    # 24 — Health Bench & Extra Credit
    f.add_table_slide(
        "Health Bench & Extra Credit",
        ["Student", "Task", "Score"],
        [
            ["Brenda", "Health Bench", "100"],
            ["James", "Health Bench", "90"],
            ["Matias", "Health Bench", "90"],
            ["Brenda", "Extra Credit (Bronx Panel)", "90"],
        ],
        "Three of you attempted the Health Bench — that's the advanced task where you "
        "evaluated the AI's clinical reasoning against a structured rubric. Brenda scored "
        "a perfect 100. James and Matias both hit 90. And Brenda also did the extra credit, "
        "running a population-level analysis on the Bronx patient panel. Her key findings: "
        "34 out of 50 patients had flu shots on record, and 12 out of 50 had an A1c of 8 "
        "or higher — which tells you something about the diabetes burden in that simulated "
        "population. These advanced tasks showed real depth of engagement with the system.",
        section="WRAP-UP"
    )

    # 25 — Common Patterns
    f.add_two_column_slide(
        "Common Patterns Across Submissions",
        "What Went Well",
        [
            "Strong medication risk identification",
            "Good use of specific chart evidence",
            "Thoughtful RAG comparison observations",
            "Most teams caught the triple whammy",
        ],
        "Growth Areas",
        [
            "RAG uncertainty calibration — some "
            "accepted RAG ON at face value",
            "Privacy reasoning sometimes stayed "
            "general vs. case-specific",
            "Citing dates alongside values "
            "(March 1 vs. March 15 matters)",
            "Checking visit notes, not just Labs tab",
        ],
        "Let me give you the honest assessment. On the strengths side, most of you "
        "were good at identifying medication risks and citing evidence from the chart. "
        "The RAG comparison observations were thoughtful. On the growth side, some of "
        "you accepted RAG ON's output without questioning whether the retrieval might "
        "be incomplete — that's the false confidence trap. And on the privacy question, "
        "a few submissions gave general essays about HIPAA rather than case-specific "
        "reasoning about Chloe. The best submissions were specific, evidence-grounded, "
        "and self-critical.",
        section="WRAP-UP"
    )

    # 26 — Five Things to Remember
    f.add_content_slide("Five Things to Remember", [
        "1. Always read visit notes, not just structured tabs — key values hide in narrative",
        "2. RAG improves grounding but can create false confidence from incomplete retrieval",
        '3. "Minimum necessary" means justify what you include AND what you exclude',
        "4. Cite specific values, dates, and medications — not vague claims",
        "5. AI governance = knowing when to trust, when to verify, and what to withhold",
    ], "Five things. Write these down. One: read the visit notes. The Labs tab doesn't "
       "have everything. Two: RAG is powerful but not perfect — don't let confident "
       "language replace your own verification. Three: minimum necessary is a two-way "
       "justification — you need to explain both what you share and what you don't. "
       "Four: be specific. 'The patient has elevated potassium' is not the same as "
       "'K+ was 5.8 mEq/L on March 12 per outside labs.' Five: AI governance isn't "
       "just about whether the AI is accurate. It's about when to trust it, when to "
       "double-check it, and when to limit what it sees.",
       section="WRAP-UP")

    # 27 — Debrief Discussion
    f.add_content_slide("Debrief Discussion", [
        "H:Questions for the Room",
        '1. "If you could redesign one thing about the ChatEHR interface to improve '
        'patient safety, what would it be?"',
        '2. "When should a clinician override an AI recommendation — even if the AI '
        'cites chart evidence to support it?"',
        '3. "How would you explain the minimum necessary principle to a hospital '
        'administrator who wants to give AI full access to all patient data?"',
    ], "Let's open it up. Three questions for you. First, if you could change one thing "
       "about the ChatEHR interface — maybe how it displays information, how the AI "
       "presents its answers, how the context levels work — what would you change to "
       "make it safer? Second, when should you override the AI even when it's citing "
       "real chart data? This gets at the false confidence issue. And third, imagine "
       "you're talking to a hospital administrator who says, 'Just give the AI full "
       "access to everything — it'll be more accurate.' How do you push back on that "
       "using what you learned from the Chloe case?",
       section="WRAP-UP",
       subtitle="Take 2 minutes to discuss with your neighbor, then we'll share out")

    # 28 — Closing
    f.add_closing_slide(
        "Thank You",
        "Next week: Module continues — bring your questions",
        "Great work this week. The fact that most of you scored in the high 80s to "
        "100 tells me you're getting it. But more importantly, the quality of your "
        "reasoning — the way you caught the triple whammy, questioned RAG's confidence, "
        "and thought through the privacy tradeoffs — that's what matters. "
        "See you next week."
    )


# =========================================================================
# Main
# =========================================================================

def main():
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)

    print("Generating 6 visual illustrations (PNG + SVG)...")
    visuals = generate_visuals()

    print("Building ChatEHR review deck...")
    factory = FordhamSlideFactory(str(TEMPLATE))
    build_opening(factory, visuals)
    build_q1_review(factory, visuals)
    build_q2_review(factory, visuals)
    build_q3_review(factory, visuals)
    build_wrapup(factory, visuals)
    factory.save(str(OUTPUT))

    # ---- Verification ----
    print("\n--- Verification ---")
    slide_count = len(Presentation(str(OUTPUT)).slides)

    def notes_count(path):
        prs = Presentation(str(path))
        with_notes = 0
        for slide in prs.slides:
            try:
                text = slide.notes_slide.notes_text_frame.text.strip()
                if text:
                    with_notes += 1
            except Exception:
                pass
        return with_notes

    n_notes = notes_count(OUTPUT)

    png_files = list(IMAGES_DIR.glob("*.png"))
    svg_files = list(IMAGES_DIR.glob("*.svg"))

    print(f"Review deck:    {slide_count} slides, {n_notes} with notes")
    print(f"PNG images:     {len(png_files)}")
    print(f"SVG images:     {len(svg_files)}")

    assert 25 <= slide_count <= 32, f"Deck has {slide_count} slides (expected 25-32)"
    assert n_notes == slide_count, (
        f"Notes: {n_notes}/{slide_count} slides have notes (expected all)")
    assert len(png_files) >= 6, f"Only {len(png_files)} PNG files (expected >= 6)"
    assert len(svg_files) >= 6, f"Only {len(svg_files)} SVG files (expected >= 6)"

    print("\nAll assertions passed!")
    print(f"\nDone! File saved to {OUTPUT}")


if __name__ == "__main__":
    main()
